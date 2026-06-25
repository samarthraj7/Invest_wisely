"""Step 1 - Deterministic deck ingestion (no LLM).

Parses PDF or PPTX into a normalized structure: per-page/slide text plus a
count of embedded images. Page numbers flow downstream as deck provenance.

Raises `DeckParseError` with an actionable message when a deck can't be read
(missing, unsupported, encrypted, or corrupt). Non-fatal issues (e.g. an
image-only / scanned deck with little extractable text) are surfaced as
`warnings` so the pipeline can continue and the report can flag low confidence.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

# Below this many extracted characters we treat the deck as effectively
# text-empty (likely scanned / image-only) and warn rather than fail.
_MIN_MEANINGFUL_CHARS = 40

IMAGE_ONLY_WARNING = (
    "Very little selectable text was extracted - this deck looks image-only "
    "or scanned. Analysis will rely on limited text; consider an OCR pass or a "
    "text-based export of the deck for higher confidence."
)


class DeckParseError(Exception):
    """Raised when a deck cannot be parsed at all."""


@dataclass
class ParsedSlide:
    page: int
    text: str
    image_count: int = 0
    ocr: bool = False  # True if this page's text came from an OCR pass


@dataclass
class ParsedDeck:
    filename: str
    kind: str  # "pdf" | "pptx"
    slides: list[ParsedSlide] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        return "\n\n".join(f"[Page {s.page}]\n{s.text}" for s in self.slides)

    @property
    def total_chars(self) -> int:
        return sum(len(s.text) for s in self.slides)

    @property
    def total_images(self) -> int:
        return sum(s.image_count for s in self.slides)


def parse_deck(path: str | Path) -> ParsedDeck:
    path = Path(path)
    if not path.exists():
        raise DeckParseError(f"File not found: {path.name}")
    if path.stat().st_size == 0:
        raise DeckParseError(f"'{path.name}' is empty (0 bytes).")

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        deck = _parse_pdf(path)
    elif suffix in (".pptx", ".ppt"):
        deck = _parse_pptx(path)
    else:
        raise DeckParseError(f"Unsupported deck type '{suffix}'. Upload a .pdf or .pptx file.")

    _add_quality_warnings(deck)
    return deck


def _add_quality_warnings(deck: ParsedDeck) -> None:
    if not deck.slides:
        raise DeckParseError("No pages/slides found in the deck.")
    if deck.total_chars < _MIN_MEANINGFUL_CHARS:
        deck.warnings.append(IMAGE_ONLY_WARNING)


def _parse_pdf(path: Path) -> ParsedDeck:
    try:
        import fitz  # PyMuPDF
    except Exception as exc:  # pragma: no cover
        raise DeckParseError("PDF support unavailable (PyMuPDF not installed).") from exc

    deck = ParsedDeck(filename=path.name, kind="pdf")
    try:
        doc = fitz.open(path)
    except Exception as exc:
        raise DeckParseError(f"'{path.name}' could not be opened - it may be corrupt.") from exc

    try:
        if getattr(doc, "needs_pass", False):
            raise DeckParseError(f"'{path.name}' is password-protected. Upload an unlocked copy.")
        for i, page in enumerate(doc, start=1):
            try:
                text = page.get_text("text").strip()
                images = len(page.get_images(full=True))
            except Exception:
                text, images = "", 0
                deck.warnings.append(f"Page {i} could not be fully read; skipped its content.")
            deck.slides.append(ParsedSlide(page=i, text=text, image_count=images))
    finally:
        doc.close()
    return deck


def _parse_pptx(path: Path) -> ParsedDeck:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise DeckParseError("PPTX support unavailable (python-pptx not installed).") from exc

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise DeckParseError(
            f"'{path.name}' could not be opened as PPTX - it may be corrupt or an old .ppt format. "
            "Re-save as .pptx or export to PDF."
        ) from exc

    deck = ParsedDeck(filename=path.name, kind="pptx")
    for i, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        image_count = 0
        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    chunks.append(shape.text_frame.text.strip())
                if getattr(shape, "shape_type", None) == 13:  # PICTURE
                    image_count += 1
            except Exception:
                continue
        deck.slides.append(ParsedSlide(page=i, text="\n".join(chunks), image_count=image_count))
    return deck
